from django.http import JsonResponse, HttpResponse
from django.views.decorators.csrf import csrf_exempt
from django.conf import settings
import os
import time
import json
import csv
import io
import zipfile
from pathlib import Path

# Document processing imports
from docx import Document
from docx.shared import Inches
import openpyxl
import pandas as pd
from pptx import Presentation
from pptx.util import Inches as PptxInches
import pypandoc
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
import ebooklib
from ebooklib import epub
import markdown
from odf import text, teletype
from odf.opendocument import load
import tempfile

# PDF processing
try:
    import PyPDF2
    PDF_SUPPORT = True
except ImportError:
    PDF_SUPPORT = False

# Supported formats
WORD_PROCESSING_FORMATS = ['doc', 'docx', 'rtf', 'odt', 'txt', 'md', 'json']
SPREADSHEET_FORMATS = ['xls', 'xlsx', 'ods', 'csv']
PRESENTATION_FORMATS = ['ppt', 'pptx', 'odp']
PUBLISHING_FORMATS = ['pdf', 'epub', 'html']
ALL_SUPPORTED_FORMATS = WORD_PROCESSING_FORMATS + SPREADSHEET_FORMATS + PRESENTATION_FORMATS + PUBLISHING_FORMATS

def api_info_documents(request):
    return JsonResponse({
        "message": "ConversionBird API - Document Processing",
        "version": "2.0.0",
        "endpoints": {
            "GET /api/documents/": "Get API information and available endpoints",
            "GET /api/documents/formats": "Get supported document formats",
            "POST /api/documents/convert": "Convert document to desired format",
            "POST /api/documents/merge": "Merge multiple documents into one",
            "POST /api/documents/split": "Split single document into multiple parts",
            "POST /api/documents/watermark": "Add text or image watermark to document",
            "POST /api/documents/password": "Add/remove password protection from PDF",
            "GET /uploads/documents/:filename": "Download processed document"
        },
        "supported_formats": {
            "word_processing": WORD_PROCESSING_FORMATS,
            "spreadsheet": SPREADSHEET_FORMATS,
            "presentation": PRESENTATION_FORMATS,
            "publishing": PUBLISHING_FORMATS
        },
        "features": {
            "format_conversion": "Convert between all supported document formats",
            "text_extraction": "Extract text from various document formats",
            "batch_processing": "Process multiple documents",
            "quality_preservation": "Maintain formatting and quality",
            "merge_split": "Merge multiple documents or split single documents",
            "watermarking": "Add text or image watermarks to documents",
            "password_protection": "Add/remove password protection for PDFs",
            "universal_merge": "Merge documents of any supported format into any output format"
        },
        "new_features": {
            "merge_documents": {
                "description": "Combine multiple documents into a single file - supports all formats",
                "supported_inputs": ["pdf", "docx", "xlsx", "pptx", "txt", "html", "md", "rtf", "odt", "csv"],
                "supported_outputs": ["pdf", "txt", "docx", "html", "epub"],
                "features": "Auto-converts mixed formats to target output format",
                "usage": "POST /api/documents/merge with multiple 'documents' files and 'format' parameter"
            },
            "split_documents": {
                "description": "Split single document into multiple parts",
                "split_types": ["pages", "size", "content"],
                "usage": "POST /api/documents/split with 'document' and 'split_type'"
            },
            "document_watermarking": {
                "description": "Add visible watermarks to documents",
                "watermark_types": ["text", "image"],
                "positions": ["center", "diagonal", "top-left", "top-right", "bottom-left", "bottom-right"],
                "usage": "POST /api/documents/watermark with 'watermark_text' or 'watermark_image'"
            },
            "pdf_password_protection": {
                "description": "Secure PDFs with password protection",
                "actions": ["protect", "unprotect"],
                "usage": "POST /api/documents/password with 'action' and 'password'"
            }
        }
    })

def get_document_formats(request):
    return JsonResponse({
        "supportedFormats": ALL_SUPPORTED_FORMATS,
        "categories": {
            "word_processing": WORD_PROCESSING_FORMATS,
            "spreadsheet": SPREADSHEET_FORMATS,
            "presentation": PRESENTATION_FORMATS,
            "publishing": PUBLISHING_FORMATS
        },
        "message": "Available document formats for conversion"
    })

@csrf_exempt
def convert_document(request):
    if request.method != 'POST':
        return JsonResponse({"error": "Method not allowed"}, status=405)

    if 'document' not in request.FILES:
        return JsonResponse({
            "error": "Please upload a document file",
            "supportedFormats": ALL_SUPPORTED_FORMATS
        }, status=400)

    document_file = request.FILES['document']
    format_param = request.POST.get('format')

    if not format_param:
        return JsonResponse({
            "error": "Please specify the output format",
            "supportedFormats": ALL_SUPPORTED_FORMATS
        }, status=400)

    if format_param.lower() not in ALL_SUPPORTED_FORMATS:
        return JsonResponse({
            "error": f"Conversion failed: Unsupported output format: {format_param}. Supported formats: {', '.join(ALL_SUPPORTED_FORMATS)}"
        }, status=400)

    try:
        # Get file extension
        input_filename = document_file.name
        input_format = get_file_extension(input_filename).lower()

        if not input_format:
            return JsonResponse({"error": "Unable to determine input file format"}, status=400)

        # Create temporary file
        with tempfile.NamedTemporaryFile(delete=False, suffix=f'.{input_format}') as temp_input:
            for chunk in document_file.chunks():
                temp_input.write(chunk)
            temp_input_path = temp_input.name

        try:
            # Perform conversion
            output_content, output_filename = perform_document_conversion(
                temp_input_path, input_format, format_param.lower()
            )

            # Save to media directory
            timestamp = int(time.time() * 1000)
            final_filename = f"{timestamp}-converted.{format_param.lower()}"
            output_path = os.path.join(settings.MEDIA_ROOT, 'documents', final_filename)
            os.makedirs(os.path.dirname(output_path), exist_ok=True)

            with open(output_path, 'wb') as f:
                f.write(output_content)

            return JsonResponse({
                "message": "Document conversion successful",
                "inputFormat": input_format,
                "outputFormat": format_param.lower(),
                "downloadUrl": f"/uploads/documents/{final_filename}",
                "fileName": final_filename
            })

        finally:
            # Clean up temporary file
            if os.path.exists(temp_input_path):
                os.unlink(temp_input_path)

    except Exception as e:
        return JsonResponse({"error": f"Conversion failed: {str(e)}"}, status=500)

def perform_document_conversion(input_path, input_format, output_format):
    """Perform document conversion based on input and output formats"""

    # Read input file content
    with open(input_path, 'rb') as f:
        input_content = f.read()

    try:
        # Try specific format handlers first
        if input_format in WORD_PROCESSING_FORMATS and output_format in WORD_PROCESSING_FORMATS:
            return convert_word_processing(input_content, input_format, output_format, input_path)

        elif input_format in SPREADSHEET_FORMATS and output_format in SPREADSHEET_FORMATS:
            return convert_spreadsheet(input_content, input_format, output_format, input_path)

        elif input_format in PRESENTATION_FORMATS and output_format in PRESENTATION_FORMATS:
            return convert_presentation(input_content, input_format, output_format, input_path)

        elif input_format == 'pdf' and output_format in PRESENTATION_FORMATS:
            # Special handling for PDF to presentation
            return convert_pdf_to_presentation(input_content, output_format)

        elif output_format in PUBLISHING_FORMATS:
            return convert_to_publishing(input_content, input_format, output_format, input_path)

        else:
            # Universal conversion using pandoc for all format combinations
            return convert_with_pandoc_universal(input_path, input_format, output_format)

    except Exception as e:
        # If specific conversion fails, try universal pandoc conversion
        try:
            return convert_with_pandoc_universal(input_path, input_format, output_format)
        except Exception as pandoc_error:
            # If pandoc also fails, try text extraction and recreation
            try:
                return convert_via_text_extraction(input_content, input_format, output_format)
            except Exception as text_error:
                # Final fallback: create a simple presentation with error message
                if output_format in PRESENTATION_FORMATS:
                    return create_error_presentation(str(e))
                else:
                    raise ValueError(f"All conversion methods failed: {str(e)}")

def convert_word_processing(content, input_format, output_format, file_path):
    """Convert between word processing formats"""

    if input_format == 'docx' and output_format == 'txt':
        doc = Document(io.BytesIO(content))
        text = '\n'.join([paragraph.text for paragraph in doc.paragraphs])
        return text.encode('utf-8'), f"converted.{output_format}"

    elif input_format == 'txt' and output_format == 'docx':
        text = content.decode('utf-8')
        doc = Document()
        doc.add_paragraph(text)
        output = io.BytesIO()
        doc.save(output)
        return output.getvalue(), f"converted.{output_format}"

    elif input_format == 'md' and output_format == 'html':
        md_text = content.decode('utf-8')
        html = markdown.markdown(md_text)
        return html.encode('utf-8'), f"converted.{output_format}"

    elif input_format == 'json' and output_format == 'txt':
        try:
            text_content = content.decode('utf-8').strip()
            if not text_content:
                return b"Empty JSON file", f"converted.{output_format}"
            json_data = json.loads(text_content)
            text = json.dumps(json_data, indent=2)
            return text.encode('utf-8'), f"converted.{output_format}"
        except json.JSONDecodeError as e:
            return f"Invalid JSON format: {str(e)}".encode('utf-8'), f"converted.{output_format}"

    else:
        # Use pandoc for other conversions
        return convert_with_pandoc(file_path, input_format, output_format)

def convert_spreadsheet(content, input_format, output_format, file_path):
    """Convert between spreadsheet formats"""

    if input_format in ['xlsx', 'xls'] and output_format == 'csv':
        if input_format == 'xlsx':
            df = pd.read_excel(io.BytesIO(content))
        else:
            df = pd.read_excel(io.BytesIO(content), engine='xlrd')
        csv_content = df.to_csv(index=False)
        return csv_content.encode('utf-8'), f"converted.{output_format}"

    elif input_format == 'csv' and output_format in ['xlsx', 'xls']:
        df = pd.read_csv(io.StringIO(content.decode('utf-8')))
        output = io.BytesIO()
        df.to_excel(output, index=False, engine='openpyxl')
        return output.getvalue(), f"converted.{output_format}"

    else:
        # Use pandoc or other methods
        return convert_with_pandoc(file_path, input_format, output_format)

def convert_presentation(content, input_format, output_format, file_path):
    """Convert between presentation formats"""

    if input_format in ['pptx', 'ppt'] and output_format == 'pdf':
        # Convert presentation to PDF using reportlab
        prs = Presentation(io.BytesIO(content))
        output = io.BytesIO()

        doc = SimpleDocTemplate(output, pagesize=letter)
        styles = getSampleStyleSheet()
        story = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    story.append(Paragraph(shape.text, styles['Normal']))
                    story.append(Spacer(1, 12))

        doc.build(story)
        return output.getvalue(), f"converted.{output_format}"

    elif input_format == 'pdf' and output_format in ['pptx', 'ppt']:
        # Convert PDF to presentation
        return convert_pdf_to_presentation(content, output_format)

    else:
        return convert_with_pandoc(file_path, input_format, output_format)

def convert_to_publishing(content, input_format, output_format, file_path):
    """Convert to publishing formats (PDF, EPUB, HTML)"""

    if output_format == 'pdf':
        return convert_to_pdf(content, input_format, file_path)

    elif output_format == 'epub':
        return convert_to_epub(content, input_format, file_path)

    elif output_format == 'html':
        return convert_to_html(content, input_format, file_path)

    else:
        return convert_with_pandoc(file_path, input_format, output_format)

def convert_to_pdf(content, input_format, file_path):
    """Convert document to PDF"""

    output = io.BytesIO()
    doc = SimpleDocTemplate(output, pagesize=letter)
    styles = getSampleStyleSheet()
    story = []

    if input_format == 'txt':
        text = content.decode('utf-8')
        story.append(Paragraph(text, styles['Normal']))

    elif input_format == 'docx':
        docx = Document(io.BytesIO(content))
        for paragraph in docx.paragraphs:
            if paragraph.text.strip():
                story.append(Paragraph(paragraph.text, styles['Normal']))
                story.append(Spacer(1, 12))

    elif input_format == 'md':
        md_text = content.decode('utf-8')
        html = markdown.markdown(md_text)
        story.append(Paragraph(html, styles['Normal']))

    else:
        # Extract text and convert
        text = extract_text_from_document(content, input_format)
        story.append(Paragraph(text, styles['Normal']))

    doc.build(story)
    return output.getvalue(), "converted.pdf"

def convert_to_epub(content, input_format, file_path):
    """Convert document to EPUB"""

    book = epub.EpubBook()
    book.set_identifier('conversionbird-epub')
    book.set_title('Converted Document')
    book.set_language('en')

    book.add_author('ConversionBird API')

    # Extract text content
    text_content = extract_text_from_document(content, input_format)

    # Create chapter
    c1 = epub.EpubHtml(title='Document', file_name='chap_01.xhtml', lang='en')
    c1.content = f'<h1>Converted Document</h1><p>{text_content.replace(chr(10), "</p><p>")}</p>'

    book.add_item(c1)

    # Create table of contents
    book.toc = (epub.Link('chap_01.xhtml', 'Document', 'intro'),)

    # Add navigation files
    book.add_item(epub.EpubNcx())
    book.add_item(epub.EpubNav())

    # Create spine
    book.spine = ['nav', c1]

    # Write EPUB
    output = io.BytesIO()
    epub.write_epub(output, book, {})
    return output.getvalue(), "converted.epub"

def convert_to_html(content, input_format, file_path):
    """Convert document to HTML"""

    if input_format == 'md':
        md_text = content.decode('utf-8')
        html_content = markdown.markdown(md_text)
    elif input_format == 'txt':
        text = content.decode('utf-8')
        html_content = f'<html><body><pre>{text}</pre></body></html>'
    else:
        text = extract_text_from_document(content, input_format)
        html_content = f'<html><body><p>{text.replace(chr(10), "</p><p>")}</p></body></html>'

    return html_content.encode('utf-8'), "converted.html"

def convert_with_pandoc_universal(input_path, input_format, output_format):
    """Universal conversion using pandoc for all format combinations"""

    # Create temporary output file
    with tempfile.NamedTemporaryFile(delete=False, suffix=f'.{output_format}') as temp_output:
        temp_output_path = temp_output.name

    try:
        # Enhanced pandoc arguments for better conversion
        extra_args = [
            '--standalone',
            '--wrap=none',
            '--markdown-headings=atx',
            '--extract-media=.',
        ]

        # Add format-specific arguments
        if output_format in ['pdf']:
            extra_args.extend(['--pdf-engine=pdflatex'])
        elif output_format in ['docx', 'odt']:
            extra_args.extend(['--reference-doc=None'])
        elif output_format in ['html']:
            extra_args.extend(['--self-contained', '--html-q-tags'])
        elif output_format in ['pptx', 'ppt']:
            # Special handling for PowerPoint formats
            extra_args.extend(['--slide-level=2', '--reference-doc=None'])
        elif output_format in ['odp']:
            extra_args.extend(['--reference-doc=None'])

        # Use pypandoc for conversion
        pypandoc.convert_file(
            input_path,
            output_format,
            outputfile=temp_output_path,
            extra_args=extra_args
        )

        # Read the converted file
        with open(temp_output_path, 'rb') as f:
            output_content = f.read()

        return output_content, f"converted.{output_format}"

    finally:
        if os.path.exists(temp_output_path):
            os.unlink(temp_output_path)

def convert_via_text_extraction(content, input_format, output_format):
    """Fallback conversion via text extraction and recreation"""

    try:
        # Extract text from input
        text_content = extract_text_from_document(content, input_format)

        if not text_content or text_content == "Text extraction failed":
            raise ValueError("Could not extract text from input file")

        # Convert extracted text to desired format
        if output_format == 'txt':
            return text_content.encode('utf-8'), "converted.txt"

        elif output_format == 'html':
            html_content = f"""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Converted Document</title>
    <style>
        body {{ font-family: Arial, sans-serif; line-height: 1.6; margin: 40px; }}
        pre {{ white-space: pre-wrap; }}
    </style>
</head>
<body>
    <pre>{text_content}</pre>
</body>
</html>"""
            return html_content.encode('utf-8'), "converted.html"

        elif output_format == 'pdf':
            return convert_text_to_pdf(text_content)

        elif output_format == 'docx':
            return convert_text_to_docx(text_content)

        elif output_format == 'md':
            return text_content.encode('utf-8'), "converted.md"

        elif output_format == 'json':
            json_content = json.dumps({
                "content": text_content,
                "source_format": input_format,
                "converted_at": time.time()
            }, indent=2, ensure_ascii=False)
            return json_content.encode('utf-8'), "converted.json"

        else:
            # For other formats, try pandoc with text input
            return convert_text_with_pandoc(text_content, output_format)

    except Exception as e:
        raise ValueError(f"Text extraction conversion failed: {str(e)}")

def convert_text_to_pdf(text_content):
    """Convert text content to PDF"""

    output = io.BytesIO()
    doc = SimpleDocTemplate(output, pagesize=letter)
    styles = getSampleStyleSheet()
    story = []

    # Split text into paragraphs and add to PDF
    paragraphs = text_content.split('\n\n')
    for para in paragraphs:
        if para.strip():
            story.append(Paragraph(para.replace('\n', '<br/>'), styles['Normal']))
            story.append(Spacer(1, 12))

    doc.build(story)
    return output.getvalue(), "converted.pdf"

def convert_text_to_docx(text_content):
    """Convert text content to DOCX"""

    doc = Document()
    doc.add_heading('Converted Document', 0)

    # Split text into paragraphs
    paragraphs = text_content.split('\n\n')
    for para in paragraphs:
        if para.strip():
            doc.add_paragraph(para)

    output = io.BytesIO()
    doc.save(output)
    return output.getvalue(), "converted.docx"

def convert_text_with_pandoc(text_content, output_format):
    """Convert text using pandoc"""

    # Create temporary input file with markdown format (pandoc recognizes this better)
    with tempfile.NamedTemporaryFile(mode='w', delete=False, suffix='.md', encoding='utf-8') as temp_input:
        # Convert plain text to basic markdown format
        markdown_content = text_content.replace('\n\n', '\n\n').strip()
        temp_input.write(markdown_content)
        temp_input_path = temp_input.name

    try:
        return convert_with_pandoc_universal(temp_input_path, 'markdown', output_format)
    finally:
        if os.path.exists(temp_input_path):
            os.unlink(temp_input_path)

def convert_pdf_to_presentation(content, output_format):
    """Convert PDF to presentation format"""

    try:
        # Extract text from PDF
        pdf_text = extract_text_from_pdf(content)

        if not pdf_text or pdf_text.strip() == "":
            # If no text extracted, create a basic presentation
            pdf_text = "PDF content could not be extracted. This is a placeholder slide."

        # Create presentation from extracted text
        prs = Presentation()

        # Split text into slides (approximately 500 characters per slide)
        text_parts = split_text_into_slides(pdf_text, max_chars_per_slide=500)

        for i, slide_text in enumerate(text_parts):
            # Create a title slide for the first slide, content slides for others
            if i == 0:
                slide_layout = prs.slide_layouts[0]  # Title slide
                slide = prs.slides.add_slide(slide_layout)
                title = slide.shapes.title
                title.text = "Converted PDF Document"
                if slide.placeholders[1].has_text_frame:
                    slide.placeholders[1].text = slide_text[:200]  # Subtitle
            else:
                slide_layout = prs.slide_layouts[1]  # Title and content
                slide = prs.slides.add_slide(slide_layout)
                title = slide.shapes.title
                title.text = f""
                if slide.placeholders[1].has_text_frame:
                    slide.placeholders[1].text = slide_text

        # Save presentation to bytes
        output = io.BytesIO()
        prs.save(output)
        return output.getvalue(), f"converted.{output_format}"

    except Exception as e:
        # Fallback: create a simple presentation with error message
        try:
            prs = Presentation()
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            title.text = "Conversion Error"
            if slide.placeholders[1].has_text_frame:
                slide.placeholders[1].text = f"Could not convert PDF: {str(e)}"

            output = io.BytesIO()
            prs.save(output)
            return output.getvalue(), f"converted.{output_format}"
        except Exception as fallback_error:
            raise ValueError(f"PDF to presentation conversion failed: {str(e)}")

def split_text_into_slides(text, max_chars_per_slide=500):
    """Split text into slides based on character count"""

    if len(text) <= max_chars_per_slide:
        return [text]

    slides = []
    words = text.split()
    current_slide = ""
    current_length = 0

    for word in words:
        if current_length + len(word) + 1 <= max_chars_per_slide:
            if current_slide:
                current_slide += " " + word
            else:
                current_slide = word
            current_length += len(word) + 1
        else:
            if current_slide:
                slides.append(current_slide)
            current_slide = word
            current_length = len(word)

    if current_slide:
        slides.append(current_slide)

    return slides if slides else [text]

def extract_text_from_pdf(content):
    """Extract text from PDF using PyPDF2"""

    if not PDF_SUPPORT:
        return "PDF text extraction not available. Please install PyPDF2."

    try:
        pdf_reader = PyPDF2.PdfReader(io.BytesIO(content))
        text_content = []

        for page in pdf_reader.pages:
            page_text = page.extract_text()
            if page_text.strip():
                text_content.append(page_text)

        return '\n\n'.join(text_content)

    except Exception as e:
        return f"PDF text extraction failed: {str(e)}"

def create_error_presentation(error_message):
    """Create a simple presentation with error message"""

    try:
        prs = Presentation()
        slide_layout = prs.slide_layouts[0]  # Title slide
        slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        title.text = "Conversion Error"

        if slide.placeholders[1].has_text_frame:
            slide.placeholders[1].text = f"An error occurred during conversion:\n\n{error_message}"

        output = io.BytesIO()
        prs.save(output)
        return output.getvalue(), "converted.pptx"

    except Exception as e:
        raise ValueError(f"Could not create error presentation: {str(e)}")

def convert_with_pandoc(input_path, input_format, output_format):
    """Legacy pandoc function - now uses universal version"""
    return convert_with_pandoc_universal(input_path, input_format, output_format)

def extract_text_from_document(content, input_format):
    """Extract text from various document formats with multilingual support"""

    try:
        if input_format == 'docx':
            doc = Document(io.BytesIO(content))
            text_parts = []
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_parts.append(paragraph.text)
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            text_parts.append(cell.text)
            return '\n\n'.join(text_parts)

        elif input_format in ['xlsx', 'xls']:
            try:
                if input_format == 'xlsx':
                    df = pd.read_excel(io.BytesIO(content), engine='openpyxl')
                else:
                    df = pd.read_excel(io.BytesIO(content), engine='xlrd')

                # Convert all data to strings and handle NaN values
                df = df.fillna('')
                df = df.astype(str)

                # Create formatted text output
                text_output = []
                text_output.append("Data from spreadsheet:")
                text_output.append("=" * 50)

                for index, row in df.iterrows():
                    row_text = f"Row {index + 1}: " + " | ".join([f"{col}: {val}" for col, val in row.items()])
                    text_output.append(row_text)

                return '\n'.join(text_output)

            except Exception as e:
                return f"Spreadsheet processing failed: {str(e)}"

        elif input_format == 'csv':
            try:
                # Try different encodings for better multilingual support
                encodings_to_try = ['utf-8', 'utf-16', 'latin-1', 'cp1252']

                for encoding in encodings_to_try:
                    try:
                        text_content = content.decode(encoding)
                        df = pd.read_csv(io.StringIO(text_content))

                        # Convert to string representation
                        text_output = []
                        text_output.append("CSV Data:")
                        text_output.append("-" * 30)
                        text_output.append(df.to_string(index=False))

                        return '\n'.join(text_output)
                    except (UnicodeDecodeError, pd.errors.ParserError):
                        continue

                return "Could not decode CSV file with supported encodings"

            except Exception as e:
                return f"CSV processing failed: {str(e)}"

        elif input_format == 'txt':
            # Try multiple encodings for multilingual support
            encodings_to_try = ['utf-8', 'utf-16', 'utf-32', 'latin-1', 'cp1252']

            for encoding in encodings_to_try:
                try:
                    return content.decode(encoding)
                except UnicodeDecodeError:
                    continue

            # Fallback with error handling
            return content.decode('utf-8', errors='replace')

        elif input_format == 'md':
            try:
                text = content.decode('utf-8')
                # Remove markdown formatting for plain text
                import re
                # Remove headers
                text = re.sub(r'^#+\s*', '', text, flags=re.MULTILINE)
                # Remove links
                text = re.sub(r'\[([^\]]+)\]\([^\)]+\)', r'\1', text)
                # Remove emphasis
                text = re.sub(r'\*\*([^\*]+)\*\*', r'\1', text)
                text = re.sub(r'\*([^\*]+)\*', r'\1', text)
                return text
            except UnicodeDecodeError:
                return content.decode('utf-8', errors='replace')

        elif input_format == 'json':
            try:
                text_content = content.decode('utf-8')
                json_data = json.loads(text_content)

                # Pretty print JSON as text
                return json.dumps(json_data, indent=2, ensure_ascii=False)
            except (UnicodeDecodeError, json.JSONDecodeError) as e:
                return f"JSON processing failed: {str(e)}"

        elif input_format == 'pdf':
            # PDF text extraction using PyPDF2
            return extract_text_from_pdf(content)

        elif input_format in ['pptx', 'ppt']:
            try:
                prs = Presentation(io.BytesIO(content))
                text_parts = []

                for slide_number, slide in enumerate(prs.slides, 1):
                    text_parts.append(f"Slide {slide_number}:")
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            text_parts.append(shape.text)
                    text_parts.append("")  # Empty line between slides

                return '\n'.join(text_parts)
            except Exception as e:
                return f"Presentation processing failed: {str(e)}"

        else:
            # Generic text extraction with multiple encoding attempts
            encodings_to_try = ['utf-8', 'utf-16', 'utf-32', 'latin-1', 'cp1252', 'iso-8859-1']

            for encoding in encodings_to_try:
                try:
                    return content.decode(encoding)
                except UnicodeDecodeError:
                    continue

            # Final fallback
            return content.decode('utf-8', errors='replace')

    except Exception as e:
        return f"Text extraction failed for {input_format}: {str(e)}"

def get_file_extension(filename):
    """Get file extension from filename"""
    return Path(filename).suffix[1:] if '.' in filename else ''

@csrf_exempt
def merge_documents(request):
    """Merge multiple documents into a single document"""
    if request.method != 'POST':
        return JsonResponse({"error": "Method not allowed"}, status=405)

    if 'documents' not in request.FILES:
        return JsonResponse({
            "error": "Please upload multiple documents",
            "usage": "Send multiple files with key 'documents'"
        }, status=400)

    documents = request.FILES.getlist('documents')
    if len(documents) < 2:
        return JsonResponse({
            "error": "Please upload at least 2 documents to merge",
            "uploaded_count": len(documents)
        }, status=400)

    output_format = request.POST.get('format', 'pdf')

    try:
        merged_content = merge_multiple_documents(documents, output_format)

        timestamp = int(time.time() * 1000)
        output_filename = f"{timestamp}-merged.{output_format}"
        output_path = os.path.join(settings.MEDIA_ROOT, 'documents', output_filename)
        os.makedirs(os.path.dirname(output_path), exist_ok=True)

        with open(output_path, 'wb') as f:
            f.write(merged_content)

        return JsonResponse({
            "message": f"Successfully merged {len(documents)} documents",
            "documents_merged": len(documents),
            "output_format": output_format,
            "downloadUrl": f"/uploads/documents/{output_filename}",
            "fileName": output_filename
        })

    except Exception as e:
        return JsonResponse({"error": f"Merge failed: {str(e)}"}, status=500)

@csrf_exempt
def split_document(request):
    """Split a single document into multiple documents"""
    if request.method != 'POST':
        return JsonResponse({"error": "Method not allowed"}, status=405)

    if 'document' not in request.FILES:
        return JsonResponse({
            "error": "Please upload a document to split"
        }, status=400)

    document = request.FILES['document']
    split_type = request.POST.get('split_type', 'pages')  # pages, size, content
    split_value = request.POST.get('split_value', '1')  # number of pages per split, size in MB, etc.

    try:
        split_files = split_single_document(document, split_type, split_value)

        timestamp = int(time.time() * 1000)
        zip_filename = f"{timestamp}-split-documents.zip"
        zip_path = os.path.join(settings.MEDIA_ROOT, 'documents', zip_filename)
        os.makedirs(os.path.dirname(zip_path), exist_ok=True)

        with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_DEFLATED) as zip_out:
            for filename, content in split_files:
                zip_out.writestr(filename, content)

        return JsonResponse({
            "message": f"Successfully split document into {len(split_files)} parts",
            "parts_created": len(split_files),
            "split_type": split_type,
            "split_value": split_value,
            "downloadUrl": f"/uploads/documents/{zip_filename}",
            "fileName": zip_filename,
            "parts": [filename for filename, _ in split_files]
        })

    except Exception as e:
        return JsonResponse({"error": f"Split failed: {str(e)}"}, status=500)

@csrf_exempt
def add_watermark(request):
    """Add watermark to document"""
    if request.method != 'POST':
        return JsonResponse({"error": "Method not allowed"}, status=405)

    if 'document' not in request.FILES:
        return JsonResponse({
            "error": "Please upload a document"
        }, status=400)

    document = request.FILES['document']
    watermark_text = request.POST.get('watermark_text')
    watermark_image = request.FILES.get('watermark_image')
    position = request.POST.get('position', 'center')  # center, diagonal, etc.
    opacity = float(request.POST.get('opacity', '0.3'))  # 0.0 to 1.0

    if not watermark_text and not watermark_image:
        return JsonResponse({
            "error": "Please provide either watermark_text or watermark_image"
        }, status=400)

    try:
        watermarked_content = add_document_watermark(document, watermark_text, watermark_image, position, opacity)

        timestamp = int(time.time() * 1000)
        input_format = get_file_extension(document.name)
        output_filename = f"{timestamp}-watermarked.{input_format}"
        output_path = os.path.join(settings.MEDIA_ROOT, 'documents', output_filename)
        os.makedirs(os.path.dirname(output_path), exist_ok=True)

        with open(output_path, 'wb') as f:
            f.write(watermarked_content)

        return JsonResponse({
            "message": "Watermark added successfully",
            "watermark_type": "text" if watermark_text else "image",
            "position": position,
            "opacity": opacity,
            "downloadUrl": f"/uploads/documents/{output_filename}",
            "fileName": output_filename
        })

    except Exception as e:
        return JsonResponse({"error": f"Watermark failed: {str(e)}"}, status=500)

@csrf_exempt
def pdf_password_protect(request):
    """Add or remove password protection from PDF"""
    if request.method != 'POST':
        return JsonResponse({"error": "Method not allowed"}, status=405)

    if 'document' not in request.FILES:
        return JsonResponse({
            "error": "Please upload a PDF document"
        }, status=400)

    document = request.FILES['document']
    action = request.POST.get('action', 'protect')  # protect or unprotect
    password = request.POST.get('password')

    if action == 'protect' and not password:
        return JsonResponse({
            "error": "Password is required for protection"
        }, status=400)

    if action == 'unprotect' and not password:
        return JsonResponse({
            "error": "Current password is required for unprotection"
        }, status=400)

    try:
        if action == 'protect':
            protected_content = add_pdf_password(document, password)
            action_desc = "protected"
        else:
            protected_content = remove_pdf_password(document, password)
            action_desc = "unprotected"

        timestamp = int(time.time() * 1000)
        output_filename = f"{timestamp}-password-{action_desc}.pdf"
        output_path = os.path.join(settings.MEDIA_ROOT, 'documents', output_filename)
        os.makedirs(os.path.dirname(output_path), exist_ok=True)

        with open(output_path, 'wb') as f:
            f.write(protected_content)

        return JsonResponse({
            "message": f"PDF {action_desc} successfully",
            "action": action,
            "downloadUrl": f"/uploads/documents/{output_filename}",
            "fileName": output_filename
        })

    except Exception as e:
        return JsonResponse({"error": f"Password operation failed: {str(e)}"}, status=500)

def merge_multiple_documents(documents, output_format):
    """Merge multiple documents into one - supports all formats"""
    if not documents:
        raise ValueError("No documents provided for merging")

    # For PDF output, try to merge PDFs directly, convert others to PDF first
    if output_format == 'pdf':
        return merge_as_pdf(documents)

    # For other formats, convert everything to the target format
    elif output_format in ['docx', 'txt', 'html']:
        return merge_via_conversion(documents, output_format)

    else:
        # For unsupported merge formats, convert to PDF as intermediate
        pdf_content = merge_as_pdf(documents)
        # Then convert PDF to desired format
        return convert_pdf_to_format(pdf_content, output_format)

def merge_as_pdf(documents):
    """Merge multiple documents as PDF - supports all formats"""
    if not PDF_SUPPORT:
        raise ValueError("PDF support not available")

    pdf_merger = PyPDF2.PdfMerger()

    for doc in documents:
        input_format = get_file_extension(doc.name).lower()

        if input_format == 'pdf':
            # PDF can be merged directly
            pdf_merger.append(io.BytesIO(doc.read()))
        else:
            # Convert other formats to PDF first
            try:
                pdf_content = convert_to_pdf(doc.read(), input_format, None)[0]
                pdf_merger.append(io.BytesIO(pdf_content))
            except Exception as e:
                # If conversion fails, create a text PDF with error message
                error_text = f"Failed to convert {doc.name}: {str(e)}"
                error_pdf = convert_text_to_pdf(error_text)[0]
                pdf_merger.append(io.BytesIO(error_pdf))

    output = io.BytesIO()
    pdf_merger.write(output)
    pdf_merger.close()

    return output.getvalue()

def merge_via_conversion(documents, output_format):
    """Merge documents by converting to target format"""
    if output_format == 'txt':
        # Combine all documents as text
        combined_text = ""
        for i, doc in enumerate(documents, 1):
            text = extract_text_from_document(doc.read(), get_file_extension(doc.name))
            combined_text += f"\n\n=== Document {i}: {doc.name} ===\n\n{text}"

        return combined_text.encode('utf-8')

    elif output_format == 'docx':
        # Create a single DOCX with all content
        combined_text = ""
        for i, doc in enumerate(documents, 1):
            text = extract_text_from_document(doc.read(), get_file_extension(doc.name))
            combined_text += f"\n\n=== Document {i}: {doc.name} ===\n\n{text}"

        return convert_text_to_docx(combined_text)[0]

    elif output_format == 'html':
        # Create HTML with all documents
        html_parts = ["<html><head><title>Merged Documents</title></head><body>"]

        for i, doc in enumerate(documents, 1):
            text = extract_text_from_document(doc.read(), get_file_extension(doc.name))
            html_parts.append(f"<h2>Document {i}: {doc.name}</h2>")
            html_parts.append(f"<pre>{text}</pre><hr>")

        html_parts.append("</body></html>")
        return '\n'.join(html_parts).encode('utf-8')

def convert_pdf_to_format(pdf_content, output_format):
    """Convert PDF content to other formats"""
    # Create a temporary PDF file
    with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_pdf:
        temp_pdf.write(pdf_content)
        temp_pdf_path = temp_pdf.name

    try:
        # Use the existing conversion system
        return convert_with_pandoc_universal(temp_pdf_path, 'pdf', output_format)
    finally:
        if os.path.exists(temp_pdf_path):
            os.unlink(temp_pdf_path)

def split_single_document(document, split_type, split_value):
    """Split a single document into multiple parts"""
    input_format = get_file_extension(document.name)

    if input_format == 'pdf' and PDF_SUPPORT:
        return split_pdf(document, split_type, int(split_value))
    else:
        # For non-PDF documents, split by content
        text = extract_text_from_document(document.read(), input_format)
        return split_text_document(text, split_type, int(split_value), input_format)

def split_pdf(document, split_type, split_value):
    """Split PDF document"""
    pdf_reader = PyPDF2.PdfReader(io.BytesIO(document.read()))
    total_pages = len(pdf_reader.pages)

    if split_type == 'pages':
        pages_per_split = split_value
    else:
        # For other split types, default to 1 page per split
        pages_per_split = 1

    split_files = []
    pdf_writer = PyPDF2.PdfWriter()

    for i in range(0, total_pages, pages_per_split):
        pdf_writer = PyPDF2.PdfWriter()

        for j in range(i, min(i + pages_per_split, total_pages)):
            pdf_writer.add_page(pdf_reader.pages[j])

        output = io.BytesIO()
        pdf_writer.write(output)
        split_files.append((f"split_{i//pages_per_split + 1}.pdf", output.getvalue()))

    return split_files

def split_text_document(text, split_type, split_value, input_format):
    """Split text-based document"""
    if split_type == 'size':
        # Split by approximate size (characters)
        chunk_size = split_value * 1000  # Convert KB to characters
    else:
        # Split by paragraphs or lines
        chunk_size = split_value * 10  # Approximate paragraphs per split

    chunks = []
    words = text.split()
    current_chunk = []

    for word in words:
        current_chunk.append(word)
        if len(' '.join(current_chunk)) >= chunk_size:
            chunks.append(' '.join(current_chunk))
            current_chunk = []

    if current_chunk:
        chunks.append(' '.join(current_chunk))

    split_files = []
    for i, chunk in enumerate(chunks, 1):
        if input_format == 'txt':
            content = chunk.encode('utf-8')
        elif input_format == 'docx':
            content = convert_text_to_docx(chunk)[0]
        else:
            content = convert_text_to_pdf(chunk)[0]

        split_files.append((f"split_{i}.{input_format}", content))

    return split_files

def add_document_watermark(document, watermark_text, watermark_image, position, opacity):
    """Add watermark to document"""
    input_format = get_file_extension(document.name)

    if input_format == 'pdf' and PDF_SUPPORT:
        return add_pdf_watermark(document, watermark_text, watermark_image, position, opacity)
    else:
        # For non-PDF documents, add watermark during conversion to PDF
        pdf_content = convert_to_pdf(document.read(), input_format, None)[0]
        # Create a temporary PDF file for watermarking
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_pdf:
            temp_pdf.write(pdf_content)
            temp_pdf_path = temp_pdf.name

        try:
            # Watermark the PDF
            watermarked_pdf = add_pdf_watermark_from_path(temp_pdf_path, watermark_text, watermark_image, position, opacity)
            return watermarked_pdf
        finally:
            if os.path.exists(temp_pdf_path):
                os.unlink(temp_pdf_path)

def add_pdf_watermark(document, watermark_text, watermark_image, position, opacity):
    """Add watermark to PDF"""
    if not PDF_SUPPORT:
        raise ValueError("PDF support not available")

    pdf_reader = PyPDF2.PdfReader(io.BytesIO(document.read()))
    pdf_writer = PyPDF2.PdfWriter()

    # Create watermark PDF
    watermark_pdf = create_watermark_pdf(watermark_text, watermark_image, position, opacity, pdf_reader.pages[0].mediabox)

    if watermark_pdf:
        watermark_reader = PyPDF2.PdfReader(io.BytesIO(watermark_pdf))

        # Apply watermark to each page
        for page_num in range(len(pdf_reader.pages)):
            page = pdf_reader.pages[page_num]
            if page_num < len(watermark_reader.pages):
                watermark_page = watermark_reader.pages[page_num]
                page.merge_page(watermark_page)
            pdf_writer.add_page(page)
    else:
        # If watermark creation failed, return original
        for page in pdf_reader.pages:
            pdf_writer.add_page(page)

    output = io.BytesIO()
    pdf_writer.write(output)
    return output.getvalue()

def add_pdf_watermark_from_path(pdf_path, watermark_text, watermark_image, position, opacity):
    """Add watermark to PDF from file path"""
    with open(pdf_path, 'rb') as f:
        return add_pdf_watermark(type('MockFile', (), {'read': f.read, 'name': pdf_path})(), watermark_text, watermark_image, position, opacity)

def create_watermark_pdf(watermark_text, watermark_image, position, opacity, page_size):
    """Create a watermark PDF page"""
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.colors import Color

    output = io.BytesIO()
    c = canvas.Canvas(output, pagesize=letter)

    # Set watermark color with opacity
    watermark_color = Color(0, 0, 0, alpha=opacity)

    if watermark_text:
        c.setFillColor(watermark_color)
        c.setFont("Helvetica", 60)

        # Position watermark
        width, height = letter
        if position == 'center':
            c.drawCentredString(width/2, height/2, watermark_text)
        elif position == 'diagonal':
            c.saveState()
            c.translate(width/2, height/2)
            c.rotate(45)
            c.drawCentredString(0, 0, watermark_text)
            c.restoreState()
        elif position == 'top-left':
            c.drawString(50, height - 100, watermark_text)
        elif position == 'top-right':
            c.drawRightString(width - 50, height - 100, watermark_text)
        elif position == 'bottom-left':
            c.drawString(50, 100, watermark_text)
        elif position == 'bottom-right':
            c.drawRightString(width - 50, 100, watermark_text)

    c.save()
    return output.getvalue()

def add_pdf_password(document, password):
    """Add password protection to PDF"""
    if not PDF_SUPPORT:
        raise ValueError("PDF support not available")

    pdf_reader = PyPDF2.PdfReader(io.BytesIO(document.read()))
    pdf_writer = PyPDF2.PdfWriter()

    # Add all pages to writer
    for page in pdf_reader.pages:
        pdf_writer.add_page(page)

    # Encrypt with password
    pdf_writer.encrypt(password)

    output = io.BytesIO()
    pdf_writer.write(output)
    return output.getvalue()

def remove_pdf_password(document, password):
    """Remove password protection from PDF"""
    if not PDF_SUPPORT:
        raise ValueError("PDF support not available")

    pdf_reader = PyPDF2.PdfReader(io.BytesIO(document.read()))

    # Try to decrypt if encrypted
    if pdf_reader.is_encrypted:
        if not pdf_reader.decrypt(password):
            raise ValueError("Incorrect password")

    pdf_writer = PyPDF2.PdfWriter()

    # Add all pages to writer (decrypted)
    for page in pdf_reader.pages:
        pdf_writer.add_page(page)

    output = io.BytesIO()
    pdf_writer.write(output)
    return output.getvalue()

def download_document(request, filename):
    file_path = os.path.join(settings.MEDIA_ROOT, 'documents', filename)
    if not os.path.exists(file_path):
        return JsonResponse({"error": "File not found"}, status=404)

    with open(file_path, 'rb') as f:
        response = HttpResponse(f.read(), content_type='application/octet-stream')
        response['Content-Disposition'] = f'attachment; filename="{filename}"'
        return response
